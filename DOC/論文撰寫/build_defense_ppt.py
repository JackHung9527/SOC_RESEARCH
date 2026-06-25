# -*- coding: utf-8 -*-
"""口試簡報產生器 — 鋰電池 SOC 估測方法之比較與嵌入式實作 (20 頁)
內容源自 DOC/論文撰寫/ 各章 markdown。"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

HERE = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(HERE, "figures")

# ---- 配色（電資學院淺銀灰藍系）----
NAVY   = RGBColor(0x1F, 0x38, 0x64)
BLUE   = RGBColor(0x2E, 0x74, 0xB5)
LBLUE  = RGBColor(0xDE, 0xEB, 0xF7)
GREY   = RGBColor(0x59, 0x5959 & 0xFF, 0x59)
DGREY  = RGBColor(0x40, 0x40, 0x40)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GOLD   = RGBColor(0xC0, 0x8A, 0x2B)
RED    = RGBColor(0xC0, 0x3A, 0x2B)
GREEN  = RGBColor(0x2E, 0x7D, 0x32)

FONT = "Microsoft JhengHei"
FONT_EN = "Times New Roman"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

PAGE = [0]  # mutable counter


def _set(run, size=18, bold=False, color=DGREY, font=FONT, italic=False):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font


def add_text(slide, l, t, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(2)
    tf.margin_top = tf.margin_bottom = Pt(2)
    return tf


def rect(slide, l, t, w, h, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp


def para(tf, text, size=18, bold=False, color=DGREY, level=0, bullet=True,
         space_after=6, align=PP_ALIGN.LEFT, first=False, font=FONT):
    p = tf.paragraphs[0] if first and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.level = level
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(0)
    if bullet:
        prefix = "▪ " if level == 0 else "– "
        r = p.add_run(); r.text = prefix
        _set(r, size, bold, BLUE if level == 0 else GREY)
    # support inline mixed; here single run
    r = p.add_run(); r.text = text
    _set(r, size, bold, color, font=font)
    return p


def header(slide, title, sub=None):
    # 頂部色帶
    rect(slide, 0, 0, SW, Inches(1.02), NAVY)
    rect(slide, 0, Inches(1.02), SW, Inches(0.07), GOLD)
    tf = add_text(slide, Inches(0.55), Inches(0.12), Inches(11.5), Inches(0.82),
                  anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title; _set(r, 27, True, WHITE)
    if sub:
        p2 = tf.add_paragraph()
        r = p2.add_run(); r.text = sub; _set(r, 13, False, LBLUE)


def footer(slide, tag):
    PAGE[0] += 1
    tf = add_text(slide, Inches(0.4), Inches(7.04), Inches(9), Inches(0.4),
                  anchor=MSO_ANCHOR.MIDDLE)
    r = tf.paragraphs[0].add_run()
    r.text = "鋰電池 SOC 估測方法之比較與嵌入式實作　|　" + tag
    _set(r, 10, False, GREY)
    tf2 = add_text(slide, Inches(12.2), Inches(7.04), Inches(0.9), Inches(0.4),
                   anchor=MSO_ANCHOR.MIDDLE)
    p2 = tf2.paragraphs[0]; p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run(); r2.text = "%d / 20" % PAGE[0]; _set(r2, 10, False, GREY)


def new(tag, title=None, sub=None):
    s = prs.slides.add_slide(BLANK)
    if title:
        header(s, title, sub)
    footer(s, tag)
    return s


def add_image_fit(slide, path, l, t, max_w, max_h, frame=True):
    im = Image.open(path); iw, ih = im.size
    ratio = iw / ih
    w = max_w; h = Emu(int(w / ratio))
    if h > max_h:
        h = max_h; w = Emu(int(h * ratio))
    x = l + Emu(int((max_w - w) / 2))
    y = t + Emu(int((max_h - h) / 2))
    if frame:
        rect(slide, x - Pt(3), y - Pt(3), w + Pt(6), h + Pt(6), WHITE, line=BLUE)
    slide.shapes.add_picture(path, x, y, width=w, height=h)
    return x, y, w, h


def add_table(slide, data, l, t, w, h, col_w=None, fs=13, head_fs=13,
              header_fill=NAVY, zebra=True, align_first_left=True):
    rows, cols = len(data), len(data[0])
    gt = slide.shapes.add_table(rows, cols, l, t, w, h).table
    if col_w:
        tot = sum(col_w)
        for i, cw in enumerate(col_w):
            gt.columns[i].width = Emu(int(w * cw / tot))
    for r in range(rows):
        for c in range(cols):
            cell = gt.cell(r, c)
            cell.margin_left = Pt(5); cell.margin_right = Pt(4)
            cell.margin_top = Pt(2); cell.margin_bottom = Pt(2)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            txt = str(data[r][c])
            tfr = cell.text_frame; tfr.word_wrap = True
            p = tfr.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if (c == 0 and align_first_left) else PP_ALIGN.CENTER
            run = p.add_run(); run.text = txt
            if r == 0:
                _set(run, head_fs, True, WHITE)
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
            else:
                _set(run, fs, False, DGREY)
                cell.fill.solid()
                cell.fill.fore_color.rgb = LBLUE if (zebra and r % 2 == 0) else WHITE
    return gt


def caption(slide, text, l, t, w):
    tf = add_text(slide, l, t, w, Inches(0.35))
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; _set(r, 11, False, GREY, italic=True)


def sidebar_note(slide, lines, l, t, w, h, title="重點", fill=LBLUE, tcolor=NAVY):
    rect(slide, l, t, w, h, fill)
    rect(slide, l, t, Inches(0.09), h, BLUE)
    tf = add_text(slide, l + Inches(0.22), t + Inches(0.12), w - Inches(0.35), h - Inches(0.2))
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = title; _set(r, 14, True, tcolor)
    p.space_after = Pt(5)
    for ln in lines:
        pp = tf.add_paragraph(); pp.space_after = Pt(4)
        r = pp.add_run(); r.text = "• " + ln; _set(r, 13, False, DGREY)


# =====================================================================
# 1. 封面
# =====================================================================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(2.55), SW, Inches(2.5), RGBColor(0x16, 0x29, 0x4A))
rect(s, 0, Inches(2.50), SW, Inches(0.06), GOLD)
rect(s, 0, Inches(5.0), SW, Inches(0.06), GOLD)
tf = add_text(s, Inches(1.0), Inches(0.55), Inches(11.3), Inches(0.6))
r = tf.paragraphs[0].add_run(); r.text = "中原大學　電機工程學系　碩士學位論文"
_set(r, 18, False, LBLUE)
# 標題兩行各自獨立文字框、固定 Y，避免缺字型時行高壓縮造成重疊
for line_txt, ty in [("鋰電池 SOC 估測方法之", 2.72), ("比較與嵌入式實作", 3.66)]:
    tf = add_text(s, Inches(0.8), Inches(ty), Inches(11.7), Inches(0.9),
                  anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = line_txt; _set(r, 36, True, WHITE)
tf = add_text(s, Inches(0.8), Inches(4.55), Inches(11.7), Inches(0.45))
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "A Comparative Study and Embedded Implementation of SOC Estimation Methods for Li-ion Batteries"
_set(r, 15, False, LBLUE, font=FONT_EN, italic=True)
tf = add_text(s, Inches(0.8), Inches(5.45), Inches(11.7), Inches(1.5), anchor=MSO_ANCHOR.TOP)
for label, val in [("研究生　：", "[ 研究生姓名 ]"),
                   ("指導教授：", "[ 指導教授 ] 博士"),
                   ("中華民國　115 年", "")]:
    p = tf.add_paragraph(); p.alignment = PP_ALIGN.CENTER; p.space_after = Pt(6)
    r = p.add_run(); r.text = label; _set(r, 18, True, WHITE)
    if val:
        r = p.add_run(); r.text = val; _set(r, 18, False, LBLUE)

# =====================================================================
# 2. 大綱
# =====================================================================
s = new("簡報大綱", "簡報大綱")
items = [
    ("1", "研究背景與動機", "BMS、SOC 角色與兩個研究缺口"),
    ("2", "研究目的與貢獻", "同電池・同協定・同 MCU 的公平比較"),
    ("3", "文獻回顧", "等效電路模型與 SOC 方法分類"),
    ("4", "測試平台與量測校正", "系統架構、韌體骨架、INA226 校正、測試協定"),
    ("5", "三種 SOC 估測方法", "庫倫計數、EKF、動態阻抗 — 原理與實作"),
    ("6", "實測結果與方法比較", "rate-capability、動態阻抗擬合、精度/資源比較"),
    ("7", "系統整合與長期再現性", "三法並行、老化與噪聲分離"),
    ("8", "結論、決策表與未來工作", "方法選用建議、SOH 延伸"),
]
y = 1.45
for num, t1, t2 in items:
    cy = Inches(y)
    circ = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.85), cy, Inches(0.52), Inches(0.52))
    circ.fill.solid(); circ.fill.fore_color.rgb = BLUE; circ.line.fill.background()
    circ.shadow.inherit = False
    ctf = circ.text_frame; ctf.word_wrap = False
    pp = ctf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
    rr = pp.add_run(); rr.text = num; _set(rr, 18, True, WHITE)
    tf = add_text(s, Inches(1.65), cy - Inches(0.04), Inches(11), Inches(0.62),
                  anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = t1 + "　"; _set(r, 20, True, NAVY)
    r = p.add_run(); r.text = "　" + t2; _set(r, 14, False, GREY)
    y += 0.685

# =====================================================================
# 3. 研究背景
# =====================================================================
s = new("研究背景", "一、研究背景", "BMS 與 SOC 在鋰電池系統中的關鍵角色")
tf = add_text(s, Inches(0.55), Inches(1.35), Inches(7.4), Inches(5.2))
para(tf, "鋰離子電池已成為電動車、可攜式裝置與儲能系統的主力儲能元件", 18, True, NAVY, first=True, space_after=4)
para(tf, "電池管理系統（BMS）負責電壓/電流/溫度監控、均衡與保護", 16, space_after=4)
para(tf, "SOC（荷電狀態）＝可用電量 ÷ 額定容量，相當於電池的「油量表」", 16, space_after=10)
para(tf, "SOC 估測的根本困難：", 17, True, NAVY, bullet=False, space_after=4)
para(tf, "SOC 無法直接量測，只能由端電壓、電流、溫度間接推估", 15, level=1, space_after=3)
para(tf, "電池為高度非線性、時變、具極化與遲滯之電化學系統", 15, level=1, space_after=3)
para(tf, "同一端電壓在不同電流/溫度/老化下對應截然不同的真實 SOC", 15, level=1, space_after=3)
para(tf, "估測失準 → 續航誤判、容量浪費、誤觸或喪失保護", 15, level=1, color=RED, space_after=3)
sidebar_note(s, [
    "端電壓被歐姆壓降與極化「污染」",
    "Vt = Voc(SOC) − I·R₀ − Vpol",
    "需藉電池模型把端電壓「還原」成 SOC",
], Inches(8.25), Inches(1.45), Inches(4.55), Inches(2.3), title="核心難點")
sidebar_note(s, [
    "本研究聚焦 SOC 估測",
    "SOH 列為未來工作（第六章）",
    "單顆 NMC 2000 mAh、室溫、不做溫補",
], Inches(8.25), Inches(4.0), Inches(4.55), Inches(2.5), title="研究範圍界定", fill=RGBColor(0xFD,0xF3,0xE0), tcolor=GOLD)

# =====================================================================
# 4. 研究動機 — 兩個缺口
# =====================================================================
s = new("研究動機", "一、研究動機", "現有的比較研究，有兩個地方說不清楚")
rect(s, Inches(0.55), Inches(1.45), Inches(6.0), Inches(4.9), LBLUE)
rect(s, Inches(0.55), Inches(1.45), Inches(6.0), Inches(0.7), BLUE)
tf = add_text(s, Inches(0.7), Inches(1.5), Inches(5.7), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
r = tf.paragraphs[0].add_run(); r.text = "缺口 ①　在電腦上比，不代表晶片上能用"; _set(r, 17, True, WHITE)
tf = add_text(s, Inches(0.8), Inches(2.35), Inches(5.5), Inches(3.9))
para(tf, "大多數研究都在電腦／MATLAB 上比，只看「準不準」", 15, first=True, space_after=7)
para(tf, "可是演算法真正要跑的地方，是很小的微控制器（MCU）", 15, space_after=7)
para(tf, "搬進 MCU 後到底佔多少記憶體、跑多久、變多不準？很少人量過", 15, color=RED, space_after=7)
para(tf, "連商用電池晶片都選「簡單夠用」的方法、而不是理論最準的 EKF——為什麼？沒人用實測講清楚", 15, space_after=4)

rect(s, Inches(6.85), Inches(1.45), Inches(5.95), Inches(4.9), RGBColor(0xFD,0xF3,0xE0))
rect(s, Inches(6.85), Inches(1.45), Inches(5.95), Inches(0.7), GOLD)
tf = add_text(s, Inches(7.0), Inches(1.5), Inches(5.7), Inches(0.6), anchor=MSO_ANCHOR.MIDDLE)
r = tf.paragraphs[0].add_run(); r.text = "缺口 ②　大家比的基準不一樣，比了不算數"; _set(r, 17, True, WHITE)
tf = add_text(s, Inches(7.1), Inches(2.35), Inches(5.5), Inches(3.9))
para(tf, "各研究用的電池、測法、硬體都不同，數字沒辦法互相比", 15, first=True, space_after=7)
para(tf, "動態阻抗法說自己「不用初始值、又快、又適合小晶片」", 15, space_after=7)
para(tf, "但從沒人把它和 EKF 放在同一顆電池、同一套測法、同一顆晶片上公平比", 15, color=RED, space_after=7)
para(tf, "結果工程師要選方法時，找不到可以信的依據", 15, space_after=4)

# =====================================================================
# 5. 研究目的與貢獻
# =====================================================================
s = new("研究目的與貢獻", "二、研究目的與貢獻", "建立「同電池・同測試協定・同 MCU」的公平比較平台")
tf = add_text(s, Inches(0.55), Inches(1.3), Inches(12.2), Inches(0.95))
para(tf, "目的：在資源受限的嵌入式環境下，對庫倫計數、EKF、動態阻抗三種代表性方法，"
        "於同一電池、同一協定、同一 MCU 上做精度與資源的公平量化比較，並歸納方法選用建議。",
     16, True, NAVY, bullet=False, first=True, space_after=0)
cards = [
    ("貢獻 ①", "讓公平比較成立的量測方法", "不只是搭硬體，而是建立「能當真值」的量測基礎："
     "INA226 校正到 < 1‰，庫倫計數才夠格當基準真值；一套協定同時產出倍率容量與動態阻抗擾動"
     "（動態阻抗免做額外實驗）；跨輪自動化，再現性實測 CV < 0.52%", BLUE),
    ("貢獻 ②", "三方法統一嵌入式實作", "把庫倫計數、EKF、動態阻抗放在同一套韌體骨架上、"
     "只替換最核心的估測模組，讓三法在「完全相同的環境」下運行——這是公平比較的前提", GREEN),
    ("貢獻 ③", "首次同硬體公平基準＋決策表", "首次把三法放在同電池、同協定、同晶片上，"
     "不只量精度、也量資源（Flash/RAM/CPU），把文獻裡「EKF 準但重、庫倫輕但會漂」這種"
     "定性印象，變成可重現的對照表與方法選用決策表 ★ 核心貢獻", GOLD),
]
x = 0.55
for tag, title, body, col in cards:
    cw = Inches(3.97)
    rect(s, Inches(x), Inches(2.45), cw, Inches(4.0), WHITE, line=col)
    rect(s, Inches(x), Inches(2.45), cw, Inches(0.95), col)
    tf = add_text(s, Inches(x+0.18), Inches(2.5), Inches(3.6), Inches(0.85), anchor=MSO_ANCHOR.MIDDLE)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = tag; _set(r, 14, True, WHITE)
    p2 = tf.add_paragraph()
    r = p2.add_run(); r.text = title; _set(r, 16, True, WHITE)
    tf = add_text(s, Inches(x+0.2), Inches(3.55), Inches(3.6), Inches(2.8))
    para(tf, body, 13.5, bullet=False, first=True, space_after=0)
    x += 4.06

# =====================================================================
# 6. 文獻回顧 — 等效電路 + 方法分類
# =====================================================================
s = new("文獻回顧", "三、文獻回顧", "等效電路模型（ECM）與 SOC 估測方法分類")
add_image_fit(s, os.path.join(FIG, "fig2-1.png"), Inches(0.5), Inches(1.4),
              Inches(4.6), Inches(4.7))
caption(s, "圖 2-1　三種等效電路模型：內阻 / 一階 RC / 二階 RC", Inches(0.5), Inches(6.15), Inches(4.6))
data = [
    ["模型", "狀態", "暫態描述", "嵌入式適用"],
    ["內阻 Rint", "1", "僅歐姆", "極佳，動態誤差大"],
    ["戴維寧 1-RC", "2", "單一時間常數", "佳，精度/成本平衡"],
    ["二階 2-RC", "3", "快慢雙常數", "中，精度高成本高"],
]
add_table(s, data, Inches(5.4), Inches(1.5), Inches(7.4), Inches(1.6),
          col_w=[2.0, 0.9, 2.0, 2.6], fs=13)
tf = add_text(s, Inches(5.4), Inches(3.35), Inches(7.45), Inches(3.4))
para(tf, "SOC 估測方法兩大類：", 16, True, NAVY, first=True, space_after=5)
para(tf, "直接計量法：庫倫計數、OCV 查表", 14, level=1, space_after=3)
para(tf, "模型基礎法：卡爾曼濾波族（KF/EKF/UKF）、動態阻抗法", 14, level=1, space_after=8)
para(tf, "本文 EKF 採一階 RC：狀態維度低、參數辨識易，"
        "於單顆 NMC、室溫下足以支撐合理端電壓重建", 14, space_after=6)
para(tf, "商用 BMS IC（bq2x、ModelGauge）多採「庫倫＋OCV 週期修正」"
        "而非 EKF —— 正是 Flash/RAM/算力/成本的工程權衡", 14, color=RED, space_after=4)

# =====================================================================
# 7. 三方法選擇依據
# =====================================================================
s = new("方法選擇依據", "三、本文方法選擇依據", "為何選庫倫計數、EKF、動態阻抗三者？")
data = [
    ["方法", "是否實作", "類別", "選擇／排除依據"],
    ["庫倫計數", "✓ 實作", "直接計量", "必備 baseline，兼任三方法比較之 ground truth（基準真值）"],
    ["OCV 查表（獨立）", "✗ 排除", "直接計量", "需長時間靜置，與 rate-capability 協定相斥；改以 GITT 表供 EKF 用"],
    ["EKF", "✓ 實作", "模型基礎", "工業事實標準；可自初值錯誤收斂；以 GITT OCV 表為觀測方程"],
    ["動態阻抗", "✓ 實作", "模型基礎", "文獻 [1] 主推；無須初值、即時可算、嵌入式友善；複用既有 dV/dI 擾動"],
]
add_table(s, data, Inches(0.55), Inches(1.5), Inches(12.25), Inches(3.0),
          col_w=[1.9, 1.2, 1.3, 6.5], fs=14)
tf = add_text(s, Inches(0.55), Inches(4.85), Inches(12.2), Inches(2.2))
para(tf, "三者恰落在「精度─資源」權衡曲線的不同位置：", 17, True, NAVY, first=True, space_after=6)
para(tf, "庫倫計數 → 基準與真值（資源下界）", 16, level=1, color=GREEN, space_after=4)
para(tf, "EKF → 模型基礎之精度上界代表（資源代價最高）", 16, level=1, color=BLUE, space_after=4)
para(tf, "動態阻抗 → 嵌入式輕量化代表（無須初值、實驗成本低）", 16, level=1, color=GOLD, space_after=4)

# =====================================================================
# 8. 系統架構
# =====================================================================
s = new("測試平台", "四、測試平台系統架構", "雙重角色：電池激勵/真值量測平台　＋　嵌入式估測標的")
add_image_fit(s, os.path.join(FIG, "fig3-1.png"), Inches(0.5), Inches(1.4),
              Inches(7.2), Inches(4.7))
caption(s, "圖 3-1　測試平台系統架構（上位機經序列埠驅動電源/負載，INA226 經 10 mΩ shunt 量測）",
        Inches(0.5), Inches(6.1), Inches(7.4))
data = [
    ["組成", "型號／規格"],
    ["微控制器", "STM32G071RB，64 MHz"],
    ["量測前端", "INA226，10 mΩ shunt，I²C"],
    ["直流電源", "ITECH IT6302（CC-CV）"],
    ["電子負載", "IT8512A+（CC 模式）"],
    ["上位機", "Python orchestrator"],
    ["受測電池", "NMC 2000 mAh，4.2/2.5 V"],
]
add_table(s, data, Inches(8.1), Inches(1.55), Inches(4.7), Inches(3.6),
          col_w=[1.5, 3.2], fs=13)
sidebar_note(s, [
    "同一顆 DUT、同一條量測迴路",
    "儀錶真值 vs MCU 估測逐點對齊",
    "Kelvin 四線、單點共地",
], Inches(8.1), Inches(5.35), Inches(4.7), Inches(1.55), title="公平比較的基礎")

# =====================================================================
# 9. 韌體骨架
# =====================================================================
s = new("韌體骨架", "四、韌體骨架（方塊圖）", "固定節拍 → once/loop 主架構 → 模組化資料流；三法只替換最核心的 SOC 估測模組")


def blkbox(slide, x, y, w, h, title, sub=None, fill=LBLUE, tcolor=NAVY,
           line=BLUE, ts=14, ss=10.5, rounded=True):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line; shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    tf = shp.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Pt(3)
    tf.margin_top = tf.margin_bottom = Pt(2)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = title; _set(r, ts, True, tcolor)
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r = p2.add_run(); r.text = sub; _set(r, ss, False, DGREY)
    return shp


def arrow(slide, x, y, w, h, down=False, color=BLUE):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW if down else MSO_SHAPE.RIGHT_ARROW,
        Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background(); shp.shadow.inherit = False
    return shp


# ── 第一層：系統節拍 ──
blkbox(s, 3.5, 1.5, 6.33, 0.62,
       "100 µs 系統節拍（硬體計時器）",
       "全系統共同時間基準；節拍只計數，累計到 1 s 觸發週期事件", fill=NAVY,
       tcolor=WHITE, line=GOLD, ts=15, ss=10.5)
arrow(s, 6.5, 2.18, 0.34, 0.42, down=True, color=GOLD)

# ── 第二層：once / loop 主架構 ──
blkbox(s, 2.55, 2.68, 3.6, 0.92, "once()　開機初始化一次",
       "通訊埠 / INA226 / 校正資料載入", fill=LBLUE, ts=14)
arrow(s, 6.25, 2.95, 0.62, 0.38, color=BLUE)
blkbox(s, 7.15, 2.68, 3.6, 0.92, "loop()　主迴圈反覆執行",
       "依節拍分派週期事件；感測器缺席仍續行", fill=LBLUE, ts=14)
arrow(s, 6.5, 3.66, 0.34, 0.42, down=True, color=BLUE)

# ── 第三層：模組化資料流 pipeline ──
tf = add_text(s, Inches(0.55), Inches(4.06), Inches(12.25), Inches(0.3))
p = tf.paragraphs[0]
r = p.add_run(); r.text = "模組化資料流（task / flow / process 三層狀態機，低耦合、可獨立替換）"
_set(r, 12, True, GREY)
pipe = [
    ("INA226", "量測前端", RGBColor(0xEC, 0xEC, 0xEC), NAVY),
    ("量測匯流排", "I²C 讀寫", LBLUE, NAVY),
    ("電池量測", "V / I / P", LBLUE, NAVY),
    ("校正套用", "14 點 LUT", LBLUE, NAVY),
    ("SOC 估測 ★", "三法共同掛載", RGBColor(0xE6, 0xF4, 0xEA), GREEN),
    ("每秒狀態回報", "→ 上位機", NAVY, WHITE),
]
bx, bw, gap, by, bh = 0.62, 1.74, 0.32, 4.45, 1.05
for i, (t1, t2, fc, tc) in enumerate(pipe):
    x = bx + i * (bw + gap)
    line = GREEN if "★" in t1 else (GOLD if tc == WHITE else BLUE)
    blkbox(s, x, by, bw, bh, t1, t2, fill=fc, tcolor=tc, line=line, ts=14, ss=11)
    if i < len(pipe) - 1:
        arrow(s, x + bw + 0.01, by + bh / 2 - 0.16, gap - 0.02, 0.32, color=BLUE)

# ── 底部說明 ──
sidebar_note(s, [
    "三法共讀同一筆校正後量測（共用前端）→ 杜絕「不同方法吃到不同量測」的不公平",
    "估測核心可獨立開關替換；除錯通訊每秒輸出一行狀態供上位機即時監看",
], Inches(0.55), Inches(5.78), Inches(12.25), Inches(1.12), title="設計要點")

# =====================================================================
# 10. INA226 校正
# =====================================================================
s = new("量測校正", "四、INA226 多點線性內插校正", "電流精度是整個平台的根 — 所有 SOC 方法都以電流積分為基礎")
tf = add_text(s, Inches(0.55), Inches(1.4), Inches(5.6), Inches(5.2))
para(tf, "校正前：電流系統性偏高約 15.8%", 16, True, RED, first=True, space_after=3)
para(tf, "200 mA 讀為 230.4 mA；主因 shunt 實阻偏離標稱 10 mΩ", 14, level=1, space_after=8)
para(tf, "校正設計：分方向各 7 點、共 14 點", 16, True, NAVY, space_after=3)
para(tf, "電流點 0 / 0.05 / 0.1 / 0.5 / 1.0 / 1.5 / 2.0 A", 14, level=1, space_after=3)
para(tf, "兼顧低電流 offset 區與高電流線性區", 14, level=1, space_after=3)
para(tf, "分段線性內插（非多項式），燒入 MCU flash 開機自動載入", 14, level=1, space_after=8)
para(tf, "校正後全範圍誤差 < 0.21 mA（< 1‰）", 16, True, GREEN, space_after=3)
para(tf, "1.75 A 僅 0.012%，優於 datasheet ±0.1%（不含 shunt 容差）", 14, level=1, space_after=3)
data = [
    ["目標(A)", "真值(mA)", "原始(mA)", "預測(mA)", "誤差(mA)"],
    ["0.300", "298.31", "340.20", "298.10", "−0.21"],
    ["0.750", "748.78", "851.32", "748.56", "−0.21"],
    ["1.250", "1248.46", "1418.52", "1248.29", "−0.17"],
    ["1.750", "1748.61", "1986.30", "1748.44", "−0.16"],
]
add_table(s, data, Inches(6.4), Inches(1.7), Inches(6.4), Inches(2.5),
          col_w=[1.2, 1.3, 1.3, 1.3, 1.2], fs=13, align_first_left=False)
caption(s, "表 3-3　LUT 內插驗證（4 個未列入校正表的獨立電流點）",
        Inches(6.4), Inches(4.25), Inches(6.4))
sidebar_note(s, [
    "充放電同電流點比值一致到小數第三位 → shunt 對方向無偏",
    "高電流區比值收斂 0.880 → shunt 實阻 ≈ 11.36 mΩ（+13.6%）",
], Inches(6.4), Inches(4.75), Inches(6.4), Inches(1.9), title="校正觀察")

# =====================================================================
# 11. 跨輪測試協定
# =====================================================================
s = new("測試協定", "四、跨輪自動化測試協定", "一輪 = 四組「充電→休息→放電→休息」，約 16–18 小時")
data = [
    ["組別", "流程（充電固定 0.5C）", "放電倍率", "目的"],
    ["1", "充飽 → 休 30 分 → 放電至截止 → 休 30 分", "0.5C", "基準放電曲線"],
    ["2", "充飽 → 休 30 分 → 放電至截止 → 休 30 分", "1.0C", "倍率能力"],
    ["3", "充飽 → 休 30 分 → 放電至截止 → 休 30 分", "1.5C", "倍率能力"],
    ["4", "充飽 → 休 30 分 → 放電至截止 → 休 30 分", "2.0C", "倍率能力"],
]
add_table(s, data, Inches(0.55), Inches(1.5), Inches(12.25), Inches(2.5),
          col_w=[0.9, 5.5, 1.4, 2.2], fs=14, align_first_left=False)
tf = add_text(s, Inches(0.55), Inches(4.35), Inches(6.1), Inches(2.6))
para(tf, "CC-CV 充電終止：V ≥ Vcv 且電流衰減至 0.1C 以下", 14, True, NAVY, first=True, space_after=8)
para(tf, "放電中注入 dV/dI 擾動：", 14, True, NAVY, space_after=3)
para(tf, "每 60 s 步降至 0.2C、dwell 1 s，取前後穩態算 ΔV/ΔI", 13, level=1, space_after=3)
para(tf, "直接服務動態阻抗法；庫倫涵蓋擾動秒數，避免約 3% 低估", 13, level=1, space_after=3)
sidebar_note(s, [
    "cycle_log.csv 跨輪持久化：cycle_id / round_id / 容量保持率 / 累積 Ah",
    "跨次執行自動續接 round_id，無縫累積",
    "三輪 fresh-cell baseline 作為資料品質閘門",
], Inches(6.85), Inches(4.35), Inches(5.95), Inches(2.55), title="持久化與品質閘門")

# =====================================================================
# 12. 庫倫計數
# =====================================================================
s = new("方法一：庫倫計數", "五、庫倫計數法（兼基準真值）", "原理 → 演算法 → 嵌入式實作 → 實驗結果")
add_image_fit(s, os.path.join(FIG, "fig4-1.png"), Inches(0.5), Inches(1.4),
              Inches(3.7), Inches(4.6))
caption(s, "圖 4-1　庫倫計數 SOC 更新流程（每秒一次）", Inches(0.5), Inches(6.05), Inches(3.7))
tf = add_text(s, Inches(4.5), Inches(1.4), Inches(8.3), Inches(2.0))
para(tf, "原理：對電流積分累計電量變化", 16, True, NAVY, first=True, space_after=3)
para(tf, "SOC(t) = SOC(t₀) − (1/C_rated) ∫ I(τ) dτ　（放電為正）", 15, level=1, font=FONT_EN, space_after=8)
para(tf, "嵌入式實作：三法中最輕量 — 每秒一乘一加，無查表、無矩陣 → 資源下界", 15, space_after=3)
sidebar_note(s, [
    "它只靠電流加總算出來，不靠電池模型，不會因為模型不準而出錯",
    "電流已經校到誤差小於千分之一，加總的底子很可靠",
    "每次放電都用「這次實際放出的總電量」當刻度：充滿算 100%、放到底算 0%，頭尾都對得準",
], Inches(4.5), Inches(3.5), Inches(8.3), Inches(1.85), title="為什麼它能當「標準答案」", fill=RGBColor(0xE6,0xF4,0xEA), tcolor=GREEN)
sidebar_note(s, [
    "一開始要先知道現在大概幾 %；如果一開始就猜錯，之後會一直錯下去",
    "電流量測只要有一點點固定偏差，一直加總就會越積越大",
    "電池會老化、能裝的電變少，用固定容量去算就會偏——而且它自己不會發現、也不會修正",
], Inches(4.5), Inches(5.5), Inches(8.3), Inches(1.5), title="它單獨用會有的三個問題", fill=RGBColor(0xFB,0xEA,0xEA), tcolor=RED)

# =====================================================================
# 13. rate-capability 實測
# =====================================================================
s = new("實測：容量再現性", "六、rate-capability 實測（庫倫，實測）", "fresh-cell 三輪各倍率放電容量 — ground truth 有效性驗證")
data = [
    ["倍率", "Round 1 (mAh)", "Round 2 (mAh)", "Round 3 (mAh)", "平均保持率"],
    ["0.5C", "1677.6", "1661.9", "1657.5", "~100.0%"],
    ["1.0C", "1651.8", "1660.0", "1658.7", "~99.5%"],
    ["1.5C", "1646.7", "1657.4", "1655.7", "~99.3%"],
    ["2.0C", "1644.9", "1652.1", "1647.4", "~98.9%"],
]
add_table(s, data, Inches(0.55), Inches(1.55), Inches(7.2), Inches(2.5),
          col_w=[1.0, 1.7, 1.7, 1.7, 1.6], fs=14, align_first_left=False)
caption(s, "表 4-1　fresh-cell 三輪各倍率實測放電容量（cycles 1–13）",
        Inches(0.55), Inches(4.2), Inches(7.2))
sidebar_note(s, [
    "rate capability 極平坦：0.5C→2.0C 容量僅降約 1%",
    "→ 內阻低、極化小；高倍率端電壓仍接近 OCV，模型失配小",
    "→ 為後續模型方法（EKF/動態阻抗）比較的有利條件",
], Inches(8.0), Inches(1.55), Inches(4.8), Inches(2.6), title="觀察 ①", fill=RGBColor(0xE6,0xF4,0xEA), tcolor=GREEN)
sidebar_note(s, [
    "同倍率跨輪再現性佳（0.5C 三輪 1657–1678 mAh）",
    "佐證測試協定 + INA226 校正之穩定性",
    "此滿放容量即 SOC 真值正規化所用的 q_full",
], Inches(0.55), Inches(4.7), Inches(12.25), Inches(1.95), title="觀察 ②")

# =====================================================================
# 14. EKF
# =====================================================================
s = new("方法二：EKF", "五、擴展卡爾曼濾波器（EKF）", "融合電流積分（預測）與端電壓（更新），可自初值錯誤收斂")
add_image_fit(s, os.path.join(FIG, "fig4-2.png"), Inches(0.5), Inches(1.4),
              Inches(4.3), Inches(4.6))
caption(s, "圖 4-2　EKF 預測—更新遞迴（一階 RC，狀態 [SOC, V₁]）",
        Inches(0.5), Inches(6.05), Inches(4.6))
tf = add_text(s, Inches(5.2), Inches(1.4), Inches(7.6), Inches(3.0))
para(tf, "狀態空間（一階 RC、狀態二維）", 16, True, NAVY, first=True, space_after=3)
para(tf, "x = [SOC, V₁]ᵀ，輸入 I，輸出 Vt", 14, level=1, font=FONT_EN, space_after=3)
para(tf, "Vt = Voc(SOC) − I·R₀ − V₁，觀測非線性源自 Voc(SOC)", 14, level=1, font=FONT_EN, space_after=3)
para(tf, "OCV–SOC 由 GITT 標準協定建 pseudo-OCV 表，1% 細網格平滑使雅可比連續", 14, level=1, space_after=8)
para(tf, "嵌入式關鍵：狀態二維、觀測一維 → 增益只需一次純量除法、免矩陣求逆", 15, True, GREEN, space_after=0)
sidebar_note(s, [
    "兼具庫倫的良好動態 + OCV 的絕對校正",
    "MCU 無硬體 FPU → 軟浮點模擬，單次更新最耗時、Flash/RAM 最大",
    "4.2.5 精度為 [待測]（依賴 GITT 表 + PC 原型 + STM32 移植）",
], Inches(5.2), Inches(4.4), Inches(7.6), Inches(2.4), title="代價與現況",
   fill=RGBColor(0xFD,0xF3,0xE0), tcolor=GOLD)

# =====================================================================
# 15. 動態阻抗 原理
# =====================================================================
s = new("方法三：動態阻抗", "五、動態阻抗法", "以端電壓對電流的瞬時變化率推估 SOC — 無須初值、無須靜置")
add_image_fit(s, os.path.join(FIG, "fig4-3.png"), Inches(0.5), Inches(1.4),
              Inches(3.7), Inches(4.6))
caption(s, "圖 4-3　離線建表 + 即時估測流程", Inches(0.5), Inches(6.05), Inches(3.7))
tf = add_text(s, Inches(4.5), Inches(1.4), Inches(8.3), Inches(3.0))
para(tf, "動態阻抗 Z = ΔV / ΔI（相鄰擾動兩點）", 16, True, NAVY, first=True, font=FONT_EN, space_after=3)
para(tf, "關鍵發現：Z 與 SOC 呈拋物線　ΔV/ΔI = a·SOC² + b·SOC + c", 15, level=1, font=FONT_EN, space_after=3)
para(tf, "滿電與放空時阻抗高、SOC≈50% 最小；對稱 → 以斜率正負判分枝取唯一解", 14, level=1, space_after=8)
para(tf, "兩階段：離線建表（PC 擬合 a,b,c 存晶片）→ 即時估測（晶片偵測擾動解二次式）",
     14, space_after=6)
sidebar_note(s, [
    "複用 rate-capability 既有 dV/dI 擾動段，無須額外實驗",
    "運算僅差分 + 低階多項式 → 介於庫倫與 EKF 之間",
    "僅擾動事件（每 60 s）產生獨立估測，事件間以庫倫內插",
], Inches(4.5), Inches(4.4), Inches(8.3), Inches(2.4), title="工程優勢與特性")

# =====================================================================
# 16. 動態阻抗 實測
# =====================================================================
s = new("實測：動態阻抗", "六、動態阻抗實測結果（fresh-cell rounds 1–3）", "二次擬合 + 以擬合反推 SOC 之逐點精度")
add_image_fit(s, os.path.join(FIG, "fig4-4.png"), Inches(0.5), Inches(1.4),
              Inches(5.6), Inches(3.0))
caption(s, "圖 4-4　|ΔV/ΔI|–SOC 實測散點與二次擬合（四種放電倍率）",
        Inches(0.5), Inches(4.5), Inches(5.6))
data = [
    ["倍率", "a (mΩ)", "b (mΩ)", "c (mΩ)", "最低點", "殘差", "反推RMSE"],
    ["0.5C", "24.5", "−27.9", "66.6", "56.9%", "3.54", "9.7%"],
    ["1.0C", "14.9", "−15.4", "61.8", "51.7%", "0.62", "6.6%"],
    ["1.5C", "15.6", "−14.6", "60.3", "46.9%", "0.59", "6.6%"],
    ["2.0C", "16.2", "−14.0", "59.0", "43.3%", "0.52", "6.3%"],
    ["合併", "20.2", "−21.6", "63.6", "53.4%", "2.90", "—"],
]
add_table(s, data, Inches(6.3), Inches(1.55), Inches(6.5), Inches(2.6),
          col_w=[1.0, 1.0, 1.0, 1.0, 1.1, 0.9, 1.2], fs=12.5, align_first_left=False)
caption(s, "表 4-3　二次擬合係數與反推精度（mΩ；SOC 以分數代入）",
        Inches(6.3), Inches(4.25), Inches(6.5))
sidebar_note(s, [
    "最低點約 SOC 53%，與理論「最小值在 50% 附近」一致",
    "高倍率殘差低至 ~0.5 mΩ；低倍率 SNR 差、殘差 ~3.5 mΩ",
    "反推 RMSE 6–10%；SOC 中段(40–60%)誤差大 → 阻抗平緩 → 反推病態",
    "→ 正當化「動態阻抗離散校正 + 庫倫內插」混合策略",
], Inches(0.55), Inches(4.85), Inches(12.25), Inches(2.0), title="關鍵觀察")

# =====================================================================
# 17. 三方法比較
# =====================================================================
s = new("方法比較", "六、三方法公平比較框架", "精度 × 強健性 × 嵌入式資源（同電池・同協定・同 MCU）")
data = [
    ["方法", "精度 RMSE", "強健性（初值錯誤）", "Flash/RAM", "每次更新運算", "浮點"],
    ["庫倫計數", "[待測]", "不收斂（無修正）", "資源下界", "一乘一加", "否"],
    ["EKF", "[待測]", "可自初值錯誤收斂", "最大（另需 OCV 表)", "最高（軟浮點）", "是"],
    ["動態阻抗", "6–10%（實測）", "首個擾動即重估", "只存幾個係數", "差分+二次式", "是"],
]
add_table(s, data, Inches(0.55), Inches(1.55), Inches(12.25), Inches(2.6),
          col_w=[1.5, 1.8, 2.6, 2.3, 2.3, 0.9], fs=13.5)
caption(s, "表 4-4 / 4-6　精度與嵌入式 footprint 比較（動態阻抗為實測，餘 [待測]）",
        Inches(0.55), Inches(4.3), Inches(12.25))
sidebar_note(s, [
    "庫倫計數最省資源，但無自我修正能力",
    "EKF 精度/強健性最佳，但 Flash/RAM/運算代價最高",
    "動態阻抗居中，且實驗成本最低（複用擾動）",
], Inches(0.55), Inches(4.75), Inches(6.0), Inches(2.0), title="精度─資源權衡")
sidebar_note(s, [
    "在同一硬體實測 Flash/RAM/CPU cycles，",
    "量化「商用 IC 為何多捨 EKF」之工程權衡",
    "★ 填補文獻少見之缺口（核心貢獻）",
], Inches(6.85), Inches(4.75), Inches(5.95), Inches(2.0), title="本研究的填補",
   fill=RGBColor(0xFD,0xF3,0xE0), tcolor=GOLD)

# =====================================================================
# 18. 系統整合 + 長期再現性
# =====================================================================
s = new("系統整合", "七、系統整合與長期再現性（實測）", "三法並行架構　＋　老化訊號與量測噪聲之分離")
tf = add_text(s, Inches(0.55), Inches(1.35), Inches(5.7), Inches(2.6))
para(tf, "三法並行整合策略", 16, True, NAVY, first=True, space_after=4)
para(tf, "共用量測前端（同一筆量測，杜絕不公平）", 13, level=1, space_after=2)
para(tf, "分離估測核心（各自獨立狀態、可開關）", 13, level=1, space_after=2)
para(tf, "統一輸出 + 交叉校正掛勾", 13, level=1, space_after=2)
para(tf, "邊際成本可加性 → 並行資源 ≈ 各 footprint 之和", 13, level=1, color=GREEN, space_after=0)
data = [
    ["倍率", "平均(mAh)", "標準差", "CV"],
    ["0.5C", "1665.7", "8.64", "0.519%"],
    ["1.0C", "1656.8", "3.60", "0.217%"],
    ["1.5C", "1653.3", "4.68", "0.283%"],
    ["2.0C", "1648.1", "2.96", "0.179%"],
]
add_table(s, data, Inches(0.55), Inches(4.1), Inches(5.7), Inches(2.2),
          col_w=[1.0, 1.6, 1.3, 1.3], fs=12.5, align_first_left=False)
caption(s, "表 5-2　fresh-cell 三輪變異性（噪聲下界）", Inches(0.55), Inches(6.35), Inches(5.7))

data2 = [
    ["Round", "日期", "0.5C 保持率", "容量(mAh)"],
    ["1", "05-12", "100.76%", "1677.6"],
    ["11", "05-22", "99.30%", "1653.4"],
    ["23", "06-02", "98.27%", "1636.3"],
    ["34", "06-14", "97.45%", "1622.5"],
    ["38", "06-19", "97.03%", "1615.5"],
]
add_table(s, data2, Inches(6.55), Inches(1.5), Inches(6.25), Inches(2.55),
          col_w=[1.0, 1.3, 1.7, 1.5], fs=12.5, align_first_left=False)
caption(s, "表 5-3　0.5C 容量保持率長期趨勢（rounds 1–38）", Inches(6.55), Inches(4.1), Inches(6.25))
sidebar_note(s, [
    "三輪變異 CV < 0.52%（噪聲下界）",
    "38 輪累計衰退約 3.7%（真實老化）",
    "老化 ≈ 7× 噪聲 → 可乾淨分離",
    "→ 支撐「逐 cycle 重錨庫倫真值」之正當性，並為 SOH 延伸提供實證基礎",
], Inches(6.55), Inches(4.55), Inches(6.25), Inches(2.3), title="老化／噪聲分離",
   fill=RGBColor(0xE6,0xF4,0xEA), tcolor=GREEN)

# =====================================================================
# 19. 結論 + 決策表
# =====================================================================
s = new("結論", "八、結論與方法選用決策表", "核心命題：嵌入式資源約束下的「精度─成本權衡」")
data = [
    ["應用約束情境", "首選方法", "主要代價"],
    ["極低成本/算力（8-bit、無 FPU、Flash<8KB）", "庫倫計數 + 靜置 OCV 修正", "無自我修正、初值錯誤不收斂"],
    ["中階、需自初值錯誤恢復、長期免校正", "EKF（一階 RC）", "需 OCV 表；運算/Flash/RAM 代價最大"],
    ["工作中即時、無靜置、實驗成本敏感", "動態阻抗 + 庫倫內插", "中段反推病態；僅事件處獨立估測"],
    ["高可靠、資源充裕、追求最佳整體", "三法混合", "整合複雜度與資源佔用最高"],
]
add_table(s, data, Inches(0.55), Inches(1.5), Inches(12.25), Inches(3.0),
          col_w=[4.2, 3.3, 4.5], fs=13.5)
caption(s, "表 6-1　SOC 方法選用決策表（動態阻抗為實測，餘待 [待測] 回填升級為定量）",
        Inches(0.55), Inches(4.65), Inches(12.25))
sidebar_note(s, [
    "不存在單一最優方法 — 最優選擇取決於目標硬體資源約束與應用需求",
    "本研究價值：將此權衡由「文獻中的定性印象」轉化為「同硬體實測支撐的決策依據」",
], Inches(0.55), Inches(5.15), Inches(12.25), Inches(1.6), title="核心結論",
   fill=RGBColor(0xFD,0xF3,0xE0), tcolor=GOLD)

# =====================================================================
# 20. 限制 + 未來工作 + 致謝
# =====================================================================
s = new("未來工作", "八、研究限制與未來工作", None)
rect(s, Inches(0.55), Inches(1.4), Inches(6.0), Inches(4.7), RGBColor(0xFB,0xEA,0xEA))
rect(s, Inches(0.55), Inches(1.4), Inches(6.0), Inches(0.6), RED)
tf = add_text(s, Inches(0.75), Inches(1.45), Inches(5.6), Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
r = tf.paragraphs[0].add_run(); r.text = "研究限制"; _set(r, 18, True, WHITE)
tf = add_text(s, Inches(0.8), Inches(2.15), Inches(5.55), Inches(3.9))
for i, t in enumerate([
    "單一電池化學體系與樣本（僅一顆 NMC）",
    "室溫、無溫度補償",
    "僅做 SOC，未做 SOH（已測到可分離老化訊號）",
    "rate capability 平坦為有利條件",
    "EKF 精度與三法 footprint 部分待實測回填",
]):
    para(tf, t, 14, first=(i == 0), space_after=8)

rect(s, Inches(6.85), Inches(1.4), Inches(5.95), Inches(4.7), RGBColor(0xE6,0xF4,0xEA))
rect(s, Inches(6.85), Inches(1.4), Inches(5.95), Inches(0.6), GREEN)
tf = add_text(s, Inches(7.05), Inches(1.45), Inches(5.6), Inches(0.5), anchor=MSO_ANCHOR.MIDDLE)
r = tf.paragraphs[0].add_run(); r.text = "未來工作"; _set(r, 18, True, WHITE)
tf = add_text(s, Inches(7.1), Inches(2.15), Inches(5.5), Inches(3.9))
para(tf, "SOH 估測延伸", 15, True, NAVY, first=True, space_after=2)
para(tf, "容量衰退追蹤、內阻成長法、投影法/ICA（複用既有量測）", 13, level=1, space_after=8)
para(tf, "LFP 電池遲滯處理", 15, True, NAVY, space_after=2)
para(tf, "遲滯狀態模型擴充至 EKF；動態阻抗不依賴 OCV 之對照", 13, level=1, space_after=8)
para(tf, "資料驅動融合與雲端 fleet learning", 15, True, NAVY, space_after=2)
para(tf, "輕量 NN 融合三法輸出；多電池雲端回饋校正", 13, level=1, space_after=0)

tf = add_text(s, Inches(0.55), Inches(6.35), Inches(12.25), Inches(0.7), anchor=MSO_ANCHOR.MIDDLE)
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "敬請各位口試委員指教　—　謝謝聆聽"; _set(r, 20, True, NAVY)

out = os.path.join(HERE, "口試簡報_鋰電池SOC估測.pptx")
prs.save(out)
print("OK ->", out, "slides=", len(prs.slides._sldIdLst))
